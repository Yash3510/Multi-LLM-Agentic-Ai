import csv
import tempfile
import unittest
from pathlib import Path

from sovereign_ai.database import Database
from sovereign_ai.deliverables import DeliverableService
from sovereign_ai.sandbox import DockerSandbox
from sovereign_ai.tools import ToolRegistry


class Phase4Tests(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory(); root = Path(self.temp.name)
        self.db = Database(root / "test.db"); self.workspace = root / "workspace"; self.registry = ToolRegistry(db=self.db, workspace=self.workspace)

    def tearDown(self): self.db.close(); self.temp.cleanup()

    def test_registry_permissions_safe_paths_and_audit(self):
        names = {tool.name for tool in self.registry.list_tools()}
        self.assertTrue({"read_file", "write_file", "move_file", "copy_file", "search_files", "csv_summary", "calculate", "delete_file"} <= names)
        created = self.registry.execute_tool("write_file", {"path": "note.txt", "content": "Compressor C-101"}, permission="write")
        self.assertTrue(created["success"])
        self.assertEqual(self.registry.execute_tool("read_file", {"path": "note.txt"})["result"], "Compressor C-101")
        denied = self.registry.execute_tool("delete_file", {"path": "note.txt"}, permission="delete")
        self.assertFalse(denied["success"])
        approved = self.registry.execute_tool("execute_python", {"code": "print('safe')"}, permission="execute")
        self.assertFalse(approved["success"])
        self.assertIn("approval", approved["error"].lower())
        escaped = self.registry.execute_tool("read_file", {"path": "../secret.txt"})
        self.assertFalse(escaped["success"])
        copied = self.registry.execute_tool("copy_file", {"source": "note.txt", "destination": "copy.txt"}, permission="write")
        self.assertTrue(copied["success"])
        moved = self.registry.execute_tool("move_file", {"source": "copy.txt", "destination": "moved.txt"}, permission="write")
        self.assertTrue(moved["success"])
        escaped_artifact = self.registry.execute_tool("generate_txt", {"name": "../outside.txt", "content": "blocked"}, permission="write")
        self.assertFalse(escaped_artifact["success"])
        self.assertGreater(self.db.execute("SELECT COUNT(*) FROM audit_events WHERE action='tool_invocation'").fetchone()[0], 0)

    def test_dataset_read_calculate_and_csv_summary(self):
        source = self.workspace / "dataset.csv"; source.write_text("item,value\nA,10\nB,20\n", encoding="utf-8")
        summary = self.registry.execute_tool("csv_summary", {"path": "dataset.csv"})
        self.assertEqual(summary["result"], {"rows": 2, "columns": ["item", "value"]})
        self.assertEqual(self.registry.execute("calculate (10 + 20) * 2"), "Calculator result: 60")

    def test_real_deliverables_are_structurally_valid(self):
        service = DeliverableService(self.workspace)
        outputs = [service.text("note.txt", "local"), service.csv("data.csv", ["name", "value"], [["A", 2]])]
        outputs += [service.docx("note.docx", "Inspection Note", ["Local evidence"]), service.xlsx("data.xlsx", "Results", ["name", "value"], [["A", 2]]), service.pptx("brief.pptx", "Inspection", "Local evidence")]
        self.assertTrue(all(item["file_size"] > 0 and Path(item["file_path"]).exists() for item in outputs))
        with (self.workspace / "data.csv").open(newline="", encoding="utf-8") as stream: self.assertEqual(next(csv.reader(stream)), ["name", "value"])
        from docx import Document
        from openpyxl import load_workbook
        from pptx import Presentation
        self.assertEqual(Document(self.workspace / "note.docx").paragraphs[0].text, "Inspection Note")
        workbook = load_workbook(self.workspace / "data.xlsx", read_only=True)
        try: self.assertEqual(workbook.sheetnames, ["Results"])
        finally: workbook.close()
        self.assertEqual(Presentation(self.workspace / "brief.pptx").slides[0].shapes.title.text, "Inspection")

    def test_document_reader_tools_extract_structure(self):
        service = DeliverableService(self.workspace)
        service.docx("read.docx", "Inspection", ["Local evidence"])
        service.xlsx("read.xlsx", "Results", ["name", "value"], [["A", 2]])
        service.pptx("read.pptx", "Inspection", "Local evidence")
        self.assertEqual(self.registry.execute_tool("read_docx", {"path": "read.docx"})["result"]["paragraphs"], ["Inspection", "Local evidence"])
        self.assertEqual(self.registry.execute_tool("read_xlsx", {"path": "read.xlsx"})["result"]["sheets"]["Results"][1], ["A", 2])
        self.assertEqual(self.registry.execute_tool("read_pptx", {"path": "read.pptx"})["result"]["slides"][0]["text"], ["Inspection", "Local evidence"])

    def test_docker_sandbox_success_failure_timeout_and_network_block(self):
        sandbox = DockerSandbox(timeout=5)
        if not sandbox.available(): self.skipTest("Docker is unavailable")
        success = sandbox.run("print('SANDBOX_OK')")
        self.assertTrue(success.success, success)
        self.assertIn("SANDBOX_OK", success.stdout)
        failure = sandbox.run("raise ValueError('expected')")
        self.assertFalse(failure.success); self.assertNotEqual(failure.exit_code, 0)
        timeout = sandbox.run("while True: pass")
        self.assertTrue(timeout.timed_out)
        network = sandbox.run("import urllib.request; urllib.request.urlopen('http://example.com', timeout=1)")
        self.assertFalse(network.success)


if __name__ == "__main__": unittest.main()

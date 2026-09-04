import csv
from datetime import datetime, timezone
from pathlib import Path


class DeliverableService:
    """Generate validated local artifacts without cloud document services."""

    def __init__(self, workspace):
        self.workspace = Path(workspace).resolve(); self.workspace.mkdir(parents=True, exist_ok=True)

    def _result(self, path, task_id=None):
        return {"file_path": str(path), "file_type": path.suffix.lower().lstrip("."), "file_size": path.stat().st_size, "creation_time": datetime.now(timezone.utc).isoformat(), "task_id": task_id}

    def text(self, name, content, task_id=None):
        path = self.workspace / name; path.write_text(content, encoding="utf-8"); return self._result(path, task_id)

    def csv(self, name, headers, rows, task_id=None):
        path = self.workspace / name
        with path.open("w", newline="", encoding="utf-8") as stream:
            writer = csv.writer(stream); writer.writerow(headers); writer.writerows(rows)
        return self._result(path, task_id)

    def docx(self, name, title, paragraphs, task_id=None):
        from docx import Document
        document = Document(); document.add_heading(title, level=1)
        for paragraph in paragraphs: document.add_paragraph(paragraph)
        path = self.workspace / name; document.save(path); return self._result(path, task_id)

    def xlsx(self, name, sheet, headers, rows, task_id=None):
        from openpyxl import Workbook, load_workbook
        book = Workbook(); page = book.active; page.title = sheet; page.append(headers)
        for row in rows: page.append(row)
        path = self.workspace / name; book.save(path); load_workbook(path, read_only=True).close(); return self._result(path, task_id)

    def pptx(self, name, title, body, task_id=None):
        from pptx import Presentation
        presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title; slide.placeholders[1].text = body
        path = self.workspace / name; presentation.save(path); return self._result(path, task_id)
